from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("stage3_phase_report.pptx")

FONT = "Microsoft YaHei"
INK = RGBColor(28, 35, 43)
MUTED = RGBColor(98, 111, 126)
LINE = RGBColor(222, 227, 232)
PAPER = RGBColor(248, 250, 252)
BLUE = RGBColor(34, 118, 197)
CYAN = RGBColor(28, 151, 166)
GREEN = RGBColor(54, 143, 93)
AMBER = RGBColor(202, 132, 35)
RED = RGBColor(190, 76, 67)
PURPLE = RGBColor(114, 92, 173)


def add_textbox(slide, x, y, w, h, text, size=18, color=INK, bold=False,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.08,
                font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.word_wrap = True
    box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_multiline(slide, x, y, w, h, lines, size=16, color=INK,
                  bullet=False, spacing=0.86):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.line_spacing = spacing
        if bullet:
            p.level = 0
            p._p.get_or_add_pPr().set("marL", "228600")
            p._p.get_or_add_pPr().set("indent", "-114300")
    return box


def rect(slide, x, y, w, h, fill, line=LINE, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def card(slide, x, y, w, h, title, body, accent=BLUE):
    shape = rect(slide, x, y, w, h, RGBColor(255, 255, 255), LINE, radius=True)
    shape.adjustments[0] = 0.08
    bar = rect(slide, x, y, 0.08, h, accent, accent)
    add_textbox(slide, x + 0.18, y + 0.12, w - 0.3, 0.35, title, 14, accent, True)
    if isinstance(body, str):
        add_textbox(slide, x + 0.18, y + 0.55, w - 0.3, h - 0.68, body, 13, INK)
    else:
        add_multiline(slide, x + 0.16, y + 0.52, w - 0.28, h - 0.62, body, 12.2, INK)
    return shape


def arrow(slide, x, y, w=0.42, color=MUTED):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.28))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def header(slide, title, section="Stage3 x0-pred"):
    add_textbox(slide, 0.54, 0.22, 8.6, 0.35, section, 9.5, MUTED, False)
    add_textbox(slide, 0.52, 0.54, 9.8, 0.55, title, 24, INK, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.54), Inches(1.18), Inches(12.25), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.color.rgb = LINE
    add_textbox(slide, 10.95, 0.28, 1.8, 0.28, "wanUpsampler", 9, MUTED, False, PP_ALIGN.RIGHT)


def footer(slide, page):
    add_textbox(slide, 0.54, 7.12, 5.0, 0.22, "源码口径: current changing_resolution, stage3 compare scripts", 7.8, MUTED)
    add_textbox(slide, 12.1, 7.12, 0.6, 0.22, str(page), 8, MUTED, False, PP_ALIGN.RIGHT)


def set_background(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER


def make_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    set_background(s)
    rect(s, 0, 0, 13.333, 7.5, PAPER, PAPER)
    add_textbox(s, 0.7, 0.7, 8.8, 0.45, "Stage3 阶段汇报", 32, INK, True)
    add_textbox(s, 0.72, 1.32, 9.5, 0.42, "x0-pred 域对齐训练与 45 / 46 / 47 步模型质量比较", 18, MUTED)
    rect(s, 0.75, 2.2, 2.55, 0.72, RGBColor(239, 246, 255), RGBColor(186, 214, 245), radius=True)
    add_textbox(s, 0.95, 2.4, 2.15, 0.24, "目标: 480p -> 720p", 13, BLUE, True, PP_ALIGN.CENTER)
    rect(s, 3.55, 2.2, 2.55, 0.72, RGBColor(239, 250, 247), RGBColor(190, 225, 213), radius=True)
    add_textbox(s, 3.75, 2.4, 2.15, 0.24, "主线: x0_pred_lr", 13, GREEN, True, PP_ALIGN.CENTER)
    rect(s, 6.35, 2.2, 2.55, 0.72, RGBColor(255, 248, 235), RGBColor(236, 210, 161), radius=True)
    add_textbox(s, 6.55, 2.4, 2.15, 0.24, "证据: 5 样本视频", 13, AMBER, True, PP_ALIGN.CENTER)
    card(s, 0.75, 3.45, 3.6, 1.72, "本次汇报回答", ["Stage3 改了什么", "为什么要比较 45/46/47", "视频证据如何判断"], BLUE)
    card(s, 4.7, 3.45, 3.6, 1.72, "当前对比组", ["interp baseline", "Stage2 clean 10k", "Stage3 step45 / 46 / 47"], CYAN)
    card(s, 8.65, 3.45, 3.6, 1.72, "待补素材", ["每个样本一段五列对比视频", "插入后只更新第 5 页观察结论"], AMBER)
    footer(s, 1)

    # 2
    s = prs.slides.add_slide(blank)
    set_background(s)
    header(s, "核心变化: 从 clean LR 转到运行时 x0-pred 域")
    card(s, 0.78, 1.55, 3.85, 1.75, "Stage2 训练域", ["输入: z0_lr", "目标: z0_hr", "问题: 推理桥接处实际拿到的不是纯净 z0"], BLUE)
    arrow(s, 4.86, 2.28, 0.6, MUTED)
    card(s, 5.72, 1.55, 3.85, 1.75, "推理桥接域", ["x0_pred = sample - sigma_t * model_output", "先估计 clean latent", "再 resize 并 re-noise"], CYAN)
    arrow(s, 9.78, 2.28, 0.6, MUTED)
    card(s, 10.65, 1.55, 1.9, 1.75, "Stage3", ["输入改为", "x0_pred_lr"], GREEN)
    rect(s, 0.78, 4.05, 11.77, 1.5, RGBColor(255, 255, 255), LINE, radius=True)
    add_textbox(s, 1.08, 4.25, 11.1, 0.32, "一句话逻辑", 14, INK, True)
    add_textbox(s, 1.08, 4.72, 10.95, 0.35, "Stage3 不换模型骨架，换训练样本分布，让 resizer 学会处理推理时真实交接点的 x0_pred_lr。", 18, INK, True)
    add_textbox(s, 1.08, 5.28, 10.8, 0.28, "对应源码: lightx2v_clean_bridge.py::step_post_upsample", 10, MUTED)
    footer(s, 2)

    # 3
    s = prs.slides.add_slide(blank)
    set_background(s)
    header(s, "Stage3 数据生成链路")
    xs = [0.72, 3.1, 5.45, 7.8, 10.15]
    titles = ["源 LMDB", "加噪到 N 步", "Wan 单步 denoise", "得到 x0_pred_lr", "写入 Stage3 LMDB"]
    bodies = [
        ["z0_lr", "z0_hr", "prompt / meta"],
        ["x_t=(1-sigma)z0 + sigma noise", "N=45/46/47"],
        ["LightX2V forward", "得到 noise_pred"],
        ["x0_pred=x_t-sigma*noise_pred", "保持 LR latent 尺寸"],
        ["schema v1", "记录 denoise_step guard"],
    ]
    colors = [BLUE, AMBER, PURPLE, CYAN, GREEN]
    for i, x in enumerate(xs):
        card(s, x, 1.75, 1.8, 2.05, titles[i], bodies[i], colors[i])
        if i < 4:
            arrow(s, x + 1.94, 2.62, 0.42, MUTED)
    rect(s, 0.75, 4.55, 11.8, 1.1, RGBColor(255, 255, 255), LINE, radius=True)
    add_textbox(s, 1.02, 4.78, 10.9, 0.36, "为什么做 45 / 46 / 47 三套", 15, INK, True)
    add_textbox(s, 1.02, 5.2, 11.0, 0.26, "不同交接步的噪声水平和模型误差不同，Stage3 数据、checkpoint、推理 change_step 必须一一匹配。", 13.5, MUTED)
    footer(s, 3)

    # 4
    s = prs.slides.add_slide(blank)
    set_background(s)
    header(s, "训练口径: 输入、目标、损失拆开看")
    card(s, 0.78, 1.55, 3.1, 2.0, "模型输入", ["仅 x0_pred_lr", "shape 与 z0_lr 一致", "不把 z0_lr forward 进模型"], GREEN)
    arrow(s, 4.08, 2.42, 0.55, MUTED)
    card(s, 4.86, 1.55, 3.1, 2.0, "Stage2 架构复用", ["WanCleanLatentResizerStage2", "16 -> 256 -> 16", "scale_factor=1.5"], BLUE)
    arrow(s, 8.16, 2.42, 0.55, MUTED)
    card(s, 8.95, 1.55, 3.1, 2.0, "监督目标", ["主目标: z0_hr", "低频锚点: z0_lr", "loss = latent + low_freq + temporal"], AMBER)
    rect(s, 0.78, 4.15, 5.45, 1.5, RGBColor(255, 255, 255), LINE, radius=True)
    add_textbox(s, 1.02, 4.35, 4.9, 0.28, "当前默认训练配置", 14, INK, True)
    add_multiline(s, 1.02, 4.78, 4.9, 0.62, ["max_steps=50000, bf16, grad_accum=8", "EMA=0.9999, eval_use_ema=true, residual_skip=false"], 11.8, MUTED)
    rect(s, 6.75, 4.15, 5.3, 1.5, RGBColor(255, 255, 255), LINE, radius=True)
    add_textbox(s, 7.0, 4.35, 4.85, 0.28, "训练安全检查", 14, INK, True)
    add_multiline(s, 7.0, 4.78, 4.85, 0.62, ["读取 meta.stage3_recipe.denoise_step", "发现 LMDB 与配置不匹配时提前失败"], 11.8, MUTED)
    footer(s, 4)

    # 5
    s = prs.slides.add_slide(blank)
    set_background(s)
    header(s, "视频证据页: 5 样本五列比较")
    add_textbox(s, 0.78, 1.36, 11.4, 0.3, "每个样本视频按同一 prompt / seed 生成，横向顺序建议固定为 interp, Stage2, Stage3-45, Stage3-46, Stage3-47。", 12.5, MUTED)
    headers = [("样本视频", 0.9, 2.0, 2.15), ("重点观察", 3.35, 2.0, 4.35), ("结论", 8.0, 2.0, 3.8)]
    for text, x, y, w in headers:
        add_textbox(s, x, y, w, 0.26, text, 11, MUTED, True)
    for i in range(5):
        y = 2.42 + i * 0.78
        rect(s, 0.82, y, 2.18, 0.5, RGBColor(255, 255, 255), LINE, radius=True)
        add_textbox(s, 1.0, y + 0.13, 1.8, 0.18, f"Sample {i + 1} 视频占位", 9.8, INK, True, PP_ALIGN.CENTER)
        rect(s, 3.25, y, 4.55, 0.5, RGBColor(255, 255, 255), LINE, radius=True)
        add_textbox(s, 3.45, y + 0.12, 4.1, 0.18, "纹理清晰度 / 时序稳定 / 细节幻觉 / 运动边缘", 9.6, MUTED)
        rect(s, 8.0, y, 3.95, 0.5, RGBColor(255, 255, 255), LINE, radius=True)
        add_textbox(s, 8.2, y + 0.12, 3.5, 0.18, "待填: 最优分支与失败模式", 9.6, MUTED)
    rect(s, 0.82, 6.47, 11.13, 0.42, RGBColor(255, 248, 235), RGBColor(238, 214, 172), radius=True)
    add_textbox(s, 1.02, 6.58, 10.7, 0.18, "插入素材后只需把每行结论改成: 哪个分支最好、哪个分支有抖动/糊/假细节。", 9.5, AMBER)
    footer(s, 5)

    # 6
    s = prs.slides.add_slide(blank)
    set_background(s)
    header(s, "阶段判断与下一步")
    card(s, 0.78, 1.55, 3.55, 2.0, "已完成", ["Stage3 x0-pred 数据/训练闭环", "45/46/47 并行构建和训练脚本", "五列视频对比脚本"], GREEN)
    card(s, 4.8, 1.55, 3.55, 2.0, "待用视频验证", ["Stage3 是否稳定优于 Stage2", "45/46/47 哪个交接步更稳", "interp 是否仍有局部优势"], BLUE)
    card(s, 8.82, 1.55, 3.55, 2.0, "决策输出", ["确定默认 change_step", "保留或淘汰对应 checkpoint", "决定是否扩展更多样本训练"], AMBER)
    rect(s, 0.78, 4.25, 11.6, 1.28, RGBColor(255, 255, 255), LINE, radius=True)
    add_textbox(s, 1.05, 4.5, 11.0, 0.32, "推荐汇报结论模板", 15, INK, True)
    add_textbox(s, 1.05, 4.96, 10.9, 0.26, "若 5 样本中 Stage3-N 在细节、稳定性和伪影三项综合最优，则把 N 作为下一轮默认交接步；否则按失败模式回到数据生成分布或损失权重。", 12.5, MUTED)
    rect(s, 0.78, 6.0, 11.6, 0.5, RGBColor(239, 246, 255), RGBColor(187, 214, 243), radius=True)
    add_textbox(s, 1.05, 6.14, 10.95, 0.18, "汇报口径保持克制: 先展示视频证据，再给默认步选择，不提前承诺 Stage3 全面优于 Stage2。", 9.8, BLUE)
    footer(s, 6)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    make_ppt()
